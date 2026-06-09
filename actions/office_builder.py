"""
office_builder.py - Brahma AI office document generation

Creates PowerPoint presentations and Excel workbooks from structured inputs.
"""

from __future__ import annotations

import json
import os
import subprocess
import re
import sys
from pathlib import Path


PROJECT_NAME = "Brahma AI - Lite"
DEFAULT_OUTPUT_DIR = Path.home() / "Desktop" / "BrahmaAI"


def _sanitize_filename(name: str, default: str) -> str:
    safe = re.sub(r"[^A-Za-z0-9._ -]+", "", (name or "").strip())
    safe = re.sub(r"\s+", " ", safe).strip().replace(" ", "_")
    return safe or default


def _resolve_output_path(output_path: str | None, title: str, ext: str) -> Path:
    if output_path:
        path = Path(output_path).expanduser()
        if not path.is_absolute():
            head = path.parts[0].lower() if path.parts else ""
            tail = Path(*path.parts[1:]) if len(path.parts) > 1 else Path(path.name)
            if head in {"downloads", "download"}:
                path = Path.home() / "Downloads" / tail
            elif head == "desktop":
                path = Path.home() / "Desktop" / tail
            else:
                path = Path.cwd() / path
        if path.suffix.lower() != ext:
            path = path.with_suffix(ext)
        path.parent.mkdir(parents=True, exist_ok=True)
        return path

    DEFAULT_OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    return DEFAULT_OUTPUT_DIR / f"{_sanitize_filename(title, 'brahma_ai_output')}{ext}"


def _parse_json_arg(value, fallback):
    if value is None:
        return fallback
    if isinstance(value, (dict, list)):
        return value
    if isinstance(value, str):
        txt = value.strip()
        if not txt:
            return fallback
        try:
            return json.loads(txt)
        except Exception:
            return fallback
    return fallback


def _coerce_cell_value(value):
    if not isinstance(value, str):
        return value
    text = value.strip()
    if not text:
        return ""
    if text.startswith("="):
        return text
    if text.lower() in {"true", "false"}:
        return text.lower() == "true"
    if re.fullmatch(r"-?\d+", text):
        try:
            return int(text)
        except Exception:
            return text
    if re.fullmatch(r"-?\d+\.\d+", text):
        try:
            return float(text)
        except Exception:
            return text
    return value


def _open_file(path: Path) -> None:
    try:
        if os.name == "nt":
            os.startfile(str(path))  # type: ignore[attr-defined]
            return
        if sys.platform == "darwin":
            subprocess.Popen(["open", str(path)])
            return
        subprocess.Popen(["xdg-open", str(path)])
    except Exception:
        pass


def _import_pptx():
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt
        return Presentation, RGBColor, MSO_SHAPE, Inches, Pt
    except Exception as e:
        raise RuntimeError(
            "python-pptx is required. Install it with: pip install python-pptx"
        ) from e


def _import_openpyxl():
    try:
        from openpyxl import Workbook
        from openpyxl.chart import BarChart, LineChart, PieChart, Reference
        from openpyxl.styles import Alignment, Font, PatternFill
        from openpyxl.utils import get_column_letter
        return Workbook, BarChart, LineChart, PieChart, Reference, Alignment, Font, PatternFill, get_column_letter
    except Exception as e:
        raise RuntimeError(
            "openpyxl is required. Install it with: pip install openpyxl"
        ) from e


def _office_theme():
    return {
        "bg": "07131C",
        "panel": "0E2230",
        "accent": "00D4FF",
        "accent2": "FF8A3D",
        "text": "F1FAFF",
        "muted": "8DB7C8",
        "line": "23485E",
    }


def _add_slide_title(slide, title: str, kicker: str, rgb, pt, Inches, MSO_SHAPE):
    colors = _office_theme()
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(colors["bg"])
    bg.line.fill.background()

    marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.45), Inches(0.14), Inches(0.14))
    marker.fill.solid()
    marker.fill.fore_color.rgb = rgb(colors["accent"])
    marker.line.fill.background()

    kicker_box = slide.shapes.add_textbox(Inches(0.72), Inches(0.38), Inches(2.5), Inches(0.28))
    tf = kicker_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = (kicker or PROJECT_NAME).upper()
    r.font.name = "Aptos"
    r.font.size = pt(10)
    r.font.bold = True
    r.font.color.rgb = rgb(colors["accent"])

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.78), Inches(10.7), Inches(0.9))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = pt(28)
    r.font.bold = True
    r.font.color.rgb = rgb(colors["text"])

    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.65), Inches(3.0), Inches(0.04))
    rule.fill.solid()
    rule.fill.fore_color.rgb = rgb(colors["accent2"])
    rule.line.fill.background()


def _slides_from_outline(outline: str | None, provided: list[dict] | None) -> list[dict]:
    if provided:
        return provided

    if not outline:
        return [
            {"title": "Overview", "bullets": ["Add a slide outline or structured slides input."]},
        ]

    blocks = [b.strip() for b in re.split(r"\n\s*\n", outline.strip()) if b.strip()]
    slides = []
    for idx, block in enumerate(blocks, 1):
        lines = [ln.strip() for ln in block.splitlines() if ln.strip()]
        title = lines[0].lstrip("#").strip() if lines else f"Slide {idx}"
        bullets = []
        for ln in lines[1:]:
            bullet = re.sub(r"^[-*•\d.)\s]+", "", ln).strip()
            if bullet:
                bullets.append(bullet)
        if not bullets and len(lines) > 1:
            bullets = lines[1:]
        slides.append({"title": title, "bullets": bullets or ["Add supporting points here."]})
    return slides


def create_presentation(parameters: dict, player=None) -> str:
    Presentation, RGBColor, MSO_SHAPE, Inches, Pt = _import_pptx()
    rgb = lambda hex_str: RGBColor.from_string(hex_str)

    title = (parameters.get("title") or parameters.get("topic") or PROJECT_NAME).strip()
    subtitle = (parameters.get("subtitle") or parameters.get("audience") or "").strip()
    output_path = _resolve_output_path(parameters.get("output_path"), title, ".pptx")
    auto_open = parameters.get("auto_open", True)
    slides = _slides_from_outline(parameters.get("outline"), _parse_json_arg(parameters.get("slides"), None))
    slides = slides[:20] if slides else slides

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    colors = _office_theme()
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(colors["bg"])
    bg.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.35), Inches(0.6), Inches(3.15), Inches(6.2))
    accent.fill.solid()
    accent.fill.fore_color.rgb = rgb(colors["panel"])
    accent.line.color.rgb = rgb(colors["line"])

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(8.4), Inches(1.9))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(32)
    r.font.bold = True
    r.font.color.rgb = rgb(colors["text"])

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.72), Inches(3.0), Inches(7.9), Inches(0.8))
        tf = sub.text_frame
        tf.word_wrap = True
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = "Aptos"
        r.font.size = Pt(16)
        r.font.color.rgb = rgb(colors["muted"])

    footer = slide.shapes.add_textbox(Inches(0.72), Inches(6.6), Inches(7.2), Inches(0.3))
    tf = footer.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = f"Created by {PROJECT_NAME}"
    r.font.name = "Aptos"
    r.font.size = Pt(10)
    r.font.color.rgb = rgb(colors["accent"])

    # Content slides
    for idx, spec in enumerate(slides, 1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_slide_title(slide, spec.get("title", f"Slide {idx}"), spec.get("kicker", f"Slide {idx}"), rgb, Pt, Inches, MSO_SHAPE)

        bullets = spec.get("bullets") or []
        body = slide.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(8.0), Inches(4.7))
        tf = body.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)

        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if isinstance(bullet, dict):
                txt = bullet.get("text", "")
            else:
                txt = str(bullet)
            p.text = f"• {txt}"
            p.level = 0
            for run in p.runs:
                run.font.name = "Aptos"
                run.font.size = Pt(18)
                run.font.color.rgb = rgb(colors["text"])

        if spec.get("notes"):
            note = slide.shapes.add_textbox(Inches(0.75), Inches(6.1), Inches(11.6), Inches(0.45))
            tf = note.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = str(spec["notes"])
            r.font.name = "Aptos"
            r.font.size = Pt(10)
            r.font.color.rgb = rgb(colors["muted"])

    prs.save(output_path)
    if auto_open:
        _open_file(output_path)
    return f"Presentation created: {output_path}"


def _choose_chart(chart_type: str, WorkbookClasses):
    Workbook, BarChart, LineChart, PieChart, Reference, Alignment, Font, PatternFill, get_column_letter = WorkbookClasses
    t = (chart_type or "bar").lower()
    if t in ("line", "trend"):
        return LineChart()
    if t in ("pie", "doughnut"):
        return PieChart()
    return BarChart()


def create_spreadsheet(parameters: dict, player=None) -> str:
    WorkbookClasses = _import_openpyxl()
    Workbook, BarChart, LineChart, PieChart, Reference, Alignment, Font, PatternFill, get_column_letter = WorkbookClasses

    title = (parameters.get("title") or parameters.get("name") or "Workbook").strip()
    output_path = _resolve_output_path(parameters.get("output_path"), title, ".xlsx")
    auto_open = parameters.get("auto_open", True)
    sheets = _parse_json_arg(parameters.get("worksheets"), None) or _parse_json_arg(parameters.get("sheets"), None)
    if not sheets:
        sheets = [{
            "name": "Sheet1",
            "headers": ["Item", "Value"],
            "rows": [["Example", 1], ["Update the tool input with your real data.", ""]],
        }]

    if isinstance(sheets, dict):
        sheets = [sheets]

    wb = Workbook()
    default_sheet = wb.active
    first_sheet = True

    theme = _office_theme()
    fill_header = PatternFill("solid", fgColor=theme["accent"])
    fill_title = PatternFill("solid", fgColor=theme["bg"])
    white_font = Font(color="FFFFFF", bold=True)
    bold_font = Font(bold=True, color=theme["bg"])
    muted_font = Font(color=theme["muted"])
    center = Alignment(horizontal="center", vertical="center")
    wrap = Alignment(wrap_text=True, vertical="top")

    for sheet_spec in sheets:
        name = _sanitize_filename(sheet_spec.get("name", "Sheet"), "Sheet")[:31]
        ws = default_sheet if first_sheet else wb.create_sheet(title=name)
        ws.title = name
        first_sheet = False

        headers = sheet_spec.get("headers") or []
        rows = sheet_spec.get("rows") or []
        title_row = sheet_spec.get("title")
        start_row = 1

        if title_row:
            ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=max(len(headers), 2, len(rows[0]) if rows else 2))
            c = ws.cell(1, 1, title_row)
            c.font = Font(bold=True, size=14, color="FFFFFF")
            c.fill = fill_title
            c.alignment = Alignment(horizontal="left", vertical="center")
            start_row = 3

        if headers:
            for col_idx, header in enumerate(headers, 1):
                cell = ws.cell(start_row, col_idx, header)
                cell.fill = fill_header
                cell.font = white_font
                cell.alignment = center
            data_start = start_row + 1
        else:
            data_start = start_row

        for row_idx, row in enumerate(rows, data_start):
            for col_idx, value in enumerate(row, 1):
                cell = ws.cell(row_idx, col_idx)
                cell.value = _coerce_cell_value(value)
                cell.alignment = wrap if isinstance(value, str) and len(value) > 24 else Alignment(vertical="top")
                if row_idx % 2 == 0 and headers:
                    cell.fill = PatternFill("solid", fgColor="F4FBFF")
                if isinstance(value, (int, float)) and col_idx > 1:
                    cell.number_format = "#,##0.00"

        if headers:
            end_row = data_start + max(len(rows) - 1, 0)
            end_col = len(headers)
            ws.auto_filter.ref = f"{ws.cell(start_row,1).coordinate}:{ws.cell(start_row + max(len(rows), 1), end_col).coordinate}"
            ws.freeze_panes = ws.cell(data_start, 1)
        else:
            end_row = data_start + max(len(rows) - 1, 0)
            end_col = max((len(r) for r in rows), default=1)

        # optional chart
        chart_spec = sheet_spec.get("chart")
        if chart_spec and rows:
            chart = _choose_chart(chart_spec.get("type", "bar"), WorkbookClasses)
            chart.title = chart_spec.get("title", "")
            chart.style = 2
            chart.y_axis.title = chart_spec.get("y_axis", "")
            chart.x_axis.title = chart_spec.get("x_axis", "")

            chart_headers_row = start_row if headers else data_start - 1
            data_start_row = data_start
            data_end_row = data_start + len(rows) - 1
            if headers and len(headers) >= 2:
                data = Reference(ws, min_col=2, min_row=chart_headers_row, max_row=data_end_row)
                cats = Reference(ws, min_col=1, min_row=data_start_row, max_row=data_end_row)
                chart.add_data(data, titles_from_data=True)
                chart.set_categories(cats)
                ws.add_chart(chart, chart_spec.get("anchor", "E2"))

        # sensible widths
        for col_idx in range(1, end_col + 1):
            values = [ws.cell(r, col_idx).value for r in range(1, min(ws.max_row, 200) + 1)]
            max_len = max([len(str(v)) for v in values if v is not None] or [10])
            ws.column_dimensions[get_column_letter(col_idx)].width = min(max(max_len + 2, 12), 36)

        ws.row_dimensions[1].height = 22
        if title_row:
            ws.row_dimensions[1].height = 26

    if "Sheet" in wb.sheetnames and len(wb.sheetnames) > 1:
        wb.remove(wb["Sheet"])

    wb.save(output_path)
    if auto_open:
        _open_file(output_path)
    return f"Spreadsheet created: {output_path}"
